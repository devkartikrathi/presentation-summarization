from pptx import Presentation
from PIL import Image
import imagehash
import io
from collections import defaultdict
import os

def extract_images_from_ppt(ppt_path):
    """Extract images from PowerPoint and return their hashes."""
    prs = Presentation(ppt_path)
    image_info = []
    
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                # Convert image blob to PIL Image
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                
                # Calculate perceptual hash
                image_hash = str(imagehash.average_hash(image))
                
                image_info.append({
                    'slide_number': slide_number,
                    'shape_id': shape.shape_id,
                    'hash': image_hash,
                    'image': image
                })
    
    return image_info

def find_duplicates(image_info):
    """Find duplicate images based on their perceptual hashes."""
    hash_groups = defaultdict(list)
    
    for info in image_info:
        hash_groups[info['hash']].append(info)
    
    # Filter only groups with duplicates
    duplicates = {k: v for k, v in hash_groups.items() if len(v) > 1}
    return duplicates

def report_duplicates(duplicates):
    """Generate a report of duplicate images."""
    if not duplicates:
        return "No duplicate images found."
    
    report = []
    for img_hash, instances in duplicates.items():
        group = [f"Slide {info['slide_number']} (Shape ID: {info['shape_id']})"
                for info in instances]
        report.append(f"\nDuplicate group found in:")
        report.extend([f"  - {location}" for location in group])
    
    return "\n".join(report)

def save_unique_images(duplicates, image_info, output_dir):
    """Save one copy of each unique image."""
    os.makedirs(output_dir, exist_ok=True)
    
    # Create set of unique hashes
    unique_hashes = set()
    for info in image_info:
        if info['hash'] not in unique_hashes:
            unique_hashes.add(info['hash'])
            # Save the image
            output_path = os.path.join(output_dir, f"unique_image_{len(unique_hashes)}.png")
            info['image'].save(output_path)

def main(ppt_path, output_dir="unique_images"):
    """Main function to process PowerPoint and find duplicate images."""
    image_info = extract_images_from_ppt(ppt_path)
    duplicates = find_duplicates(image_info)
    
    # Generate and print report
    report = report_duplicates(duplicates)
    print(report)
    
    # Save unique images
    save_unique_images(duplicates, image_info, output_dir)
    print(f"\nUnique images saved to: {output_dir}")

if __name__ == "__main__":
    main("uploaded_documents/1 TB Education VDH dup.pptx")
