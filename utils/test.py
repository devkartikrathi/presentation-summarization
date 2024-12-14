import os
from dotenv import load_dotenv
import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from io import BytesIO

load_dotenv()

def process_pdf_to_presentation(pdf_file):
    api_key = os.getenv("GEMINI_API_KEY") 

    if not api_key:
        st.error("API key for Google Gemini is missing. Please set the GEMINI_API_KEY environment variable.")
        return None

    genai.configure(api_key=api_key)

    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        
        uploaded_pdf = genai.upload_file(pdf_file, mime_type="application/pdf")

        request_text = "Extract the core content of this PDF and format it as text that can be used for making a presentation."
        response = model.generate_content([request_text, uploaded_pdf])
        response = model.generate_content(f"Summarize this content to fit the same content in less number of slides without compromizing on the knowledge it provides, {response.text}")

        return response.text
    
    except Exception as e:
        st.error(f"Error calling Google Gemini API: {e}")
        
        return None
    
def create_ppt_from_content(content):
    presentation = Presentation()
    slides = content.split("\n\n")

    for slide in slides:
        lines = slide.split("\n")
        slide_layout = presentation.slide_layouts[1] 
        ppt_slide = presentation.slides.add_slide(slide_layout)
        
        if lines:
            ppt_slide.shapes.title.text = lines[0]
            if len(lines) > 1:
                ppt_slide.placeholders[1].text = "\n".join(lines[1:])  

    ppt_io = BytesIO()
    presentation.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

def main():
    st.title("PDF to Presentation Content using Google Gemini API")
    st.write("Upload a PDF file, and this app will generate content suitable for creating a presentation using the Google Gemini API.")

    uploaded_file = st.file_uploader("Upload PDF", type="pdf")

    if uploaded_file is not None:
        with st.spinner("Processing your PDF, please wait..."):
            presentation_content = process_pdf_to_presentation(uploaded_file)

        if presentation_content:
            st.subheader("Presentation Content")
            st.write(presentation_content)

            ppt_file = create_ppt_from_content(presentation_content)

            st.download_button(
                label="Download PowerPoint Presentation",
                data=ppt_file,
                file_name="summary_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        else:
            st.error("Failed to generate presentation content.")

if __name__ == "__main__":
    main()