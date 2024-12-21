import os
import base64
import io
import uuid
from typing import List, Dict, Optional, Tuple, Any
import streamlit as st
from PIL import Image
from gtts import gTTS
from PyPDF2 import PdfReader
from pptx import Presentation
import imagehash
from langchain_together import TogetherEmbeddings, ChatTogether
from langchain_chroma import Chroma
from langchain.retrievers import MultiVectorRetriever
from langchain.storage import InMemoryStore
from langchain.schema.document import Document
from langchain.prompts import ChatPromptTemplate
from langchain.schema.runnable import RunnableParallel
from langchain.schema.output_parser import StrOutputParser
from unstructured.partition.pptx import partition_pptx
from unstructured.partition.pdf import partition_pdf
from together import Together
import chromadb

# Configure page
st.set_page_config(
    page_title="Clinical Document Assistant",
    page_icon="🏥",
    layout="wide"
)

# Initialize ChromaDB
PERSIST_DIRECTORY = "chroma_db"
os.makedirs(PERSIST_DIRECTORY, exist_ok=True)

# Initialize session state
if 'processed_doc' not in st.session_state:
    st.session_state.processed_doc = None
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'rag_model' not in st.session_state:
    st.session_state.rag_model = None
if 'chroma_client' not in st.session_state:
    st.session_state.chroma_client = chromadb.PersistentClient(path=PERSIST_DIRECTORY)

class ImageProcessor:
    def __init__(self, client):
        self.client = Together()
        self.vision_model = "meta-llama/Llama-3.2-11B-Vision-Instruct-Turbo"

    def summarize_image(self, encoded_image: str) -> str:
        try:
            response = self.client.chat.completions.create(
                model=self.vision_model,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe the contents of this image in detail, focusing on medical relevance if any."}, 
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encoded_image}"}}
                    ]
                }],
                max_tokens=500,
                temperature=0.7
            )
            return response.choices[0].message.content
        except Exception as e:
            st.error(f"Error summarizing image: {str(e)}")
            return "Error processing image"

    def process_images(self, images: List[Dict]) -> List[Dict]:
        processed_images = []
        for img in images:
            if 'base64' in img:
                summary = self.summarize_image(img['base64'])
                img['summary'] = summary
                processed_images.append(img)
        return processed_images

    @staticmethod
    def _process_image(image: Image.Image) -> Image.Image:
        if image.mode in ('RGBA', 'LA'):
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[3] if image.mode == 'RGBA' else image.split()[1])
            return background
        return image.convert('RGB') if image.mode != 'RGB' else image

    @staticmethod
    def extract_images(file_path: str, is_pdf: bool = True) -> List[Dict]:
        try:
            reader = PdfReader(file_path) if is_pdf else Presentation(file_path)
            images = []
            
            if is_pdf:
                for page_num, page in enumerate(reader.pages, 1):
                    if '/XObject' in page['/Resources']:
                        for obj in page['/Resources']['/XObject'].get_object().values():
                            if obj['/Subtype'] == '/Image':
                                try:
                                    image = Image.open(io.BytesIO(obj.get_data()))
                                    processed_image = ImageProcessor._process_image(image)
                                    buffered = io.BytesIO()
                                    processed_image.save(buffered, format="JPEG", quality=85)
                                    images.append({
                                        'page': page_num,
                                        'base64': base64.b64encode(buffered.getvalue()).decode(),
                                        'metadata': {'page': page_num}
                                    })
                                except Exception as e:
                                    st.warning(f"Skipped image on page {page_num}: {str(e)}")
            else:
                for slide_num, slide in enumerate(reader.slides, 1):
                    for shape in slide.shapes:
                        if hasattr(shape, "image"):
                            try:
                                image = Image.open(io.BytesIO(shape.image.blob))
                                processed_image = ImageProcessor._process_image(image)
                                buffered = io.BytesIO()
                                processed_image.save(buffered, format="JPEG", quality=85)
                                images.append({
                                    'slide': slide_num,
                                    'base64': base64.b64encode(buffered.getvalue()).decode(),
                                    'metadata': {'slide': slide_num}
                                })
                            except Exception as e:
                                st.warning(f"Skipped image on slide {slide_num}: {str(e)}")
            return images
        except Exception as e:
            st.error(f"Error extracting images: {str(e)}")
            return []

class RAGModel:
    def __init__(self, collection_name: str = "medical_docs"):
        self.embeddings = TogetherEmbeddings(
            model="togethercomputer/m2-bert-80M-32k-retrieval"
        )
        
        try:
            self.vectorstore = Chroma(
                client=st.session_state.chroma_client,
                collection_name=collection_name,
                embedding_function=self.embeddings,
                persist_directory=PERSIST_DIRECTORY
            )
        except Exception as e:
            st.session_state.chroma_client.create_collection(collection_name)
            self.vectorstore = Chroma(
                client=st.session_state.chroma_client,
                collection_name=collection_name,
                embedding_function=self.embeddings,
                persist_directory=PERSIST_DIRECTORY
            )
        
        self.store = InMemoryStore()
        self.retriever = MultiVectorRetriever(
            vectorstore=self.vectorstore,
            docstore=self.store,
            id_key="doc_id",
            search_type="similarity",
            search_kwargs={"k": 5}
        )
        self.llm = ChatTogether(
            model="meta-llama/Llama-3.2-3B-Instruct-Turbo"
        )

    def add_documents(self, texts: List[str], images: List[Dict]) -> None:
        docs = []
        
        # Add text documents
        for text in texts:
            if text.strip(): 
                doc_id = str(uuid.uuid4())
                docs.append(Document(
                    page_content=text,
                    metadata={"doc_id": doc_id, "type": "text"}
                ))
        
        for img in images:
            if img.get('summary'):
                doc_id = str(uuid.uuid4())
                docs.append(Document(
                    page_content=img['summary'],
                    metadata={
                        "doc_id": doc_id,
                        "type": "image",
                        "base64": img.get('base64', ''),
                        **img.get('metadata', {})
                    }
                ))
        
        if docs: 
            try:
                self.vectorstore.add_documents(docs)
                self.store.mset([(doc.metadata["doc_id"], doc) for doc in docs])
            except Exception as e:
                st.error(f"Error adding documents to vectorstore: {str(e)}")

    def get_relevant_context(self, query: str) -> Tuple[List[str], List[Dict]]:
        try:
            docs = self.retriever.invoke(query)
            texts = []
            images = []
            
            for doc in docs:
                if doc.metadata.get("type") == "text":
                    texts.append(doc.page_content)
                elif doc.metadata.get("type") == "image":
                    images.append({
                        "base64": doc.metadata.get("base64", ""),
                        "metadata": {k: v for k, v in doc.metadata.items() 
                                   if k not in ["doc_id", "type", "base64"]}
                    })
            
            return texts, images
        except Exception as e:
            st.error(f"Error retrieving context: {str(e)}")
            return [], []

def process_document(file):
    if not file:
        return None
        
    try:
        temp_dir = "temp_files"
        os.makedirs(temp_dir, exist_ok=True)
        
        file_path = os.path.join(temp_dir, file.name)
        with open(file_path, "wb") as f:
            f.write(file.getvalue())
            
        is_pdf = file.name.lower().endswith('.pdf')
        
        elements = (partition_pdf if is_pdf else partition_pptx)(
            filename=file_path,
            extract_images_in_pdf=True,
            infer_table_structure=True
        )
        
        image_processor = ImageProcessor(Together())
        images = ImageProcessor.extract_images(file_path, is_pdf)
        processed_images = image_processor.process_images(images)
        
        text_content = [
            element.text for element in elements 
            if hasattr(element, 'text') and element.text.strip()
        ]
        
        if not st.session_state.rag_model:
            st.session_state.rag_model = RAGModel()
        
        st.session_state.rag_model.add_documents(text_content, processed_images)
        
        return {
            'text_elements': text_content,
            'images': processed_images,
            'doc_id': str(uuid.uuid4())
        }
        
    except Exception as e:
        st.error(f"Error processing document: {str(e)}")
        return None
    finally:
        if 'file_path' in locals():
            try:
                os.remove(file_path)
            except:
                pass

def main():
    # Sidebar
    with st.sidebar:
        st.title("Document Upload")
        uploaded_file = st.file_uploader(
            "Upload Medical Document",
            type=['pdf', 'pptx'],
            help="Upload a PDF or PowerPoint file"
        )
        
        if uploaded_file and not st.session_state.processed_doc:
            with st.spinner("Processing document..."):
                st.session_state.processed_doc = process_document(uploaded_file)
                if st.session_state.processed_doc:
                    st.success("Document processed successfully!")
    
    # Main content
    st.title("Clinical Document Assistant 🏥")
    
    if st.session_state.processed_doc:
        tab1, tab2, tab3, tab4 = st.tabs(["Chat", "Summary", "Slides", "Images"])
        
        with tab1:
            st.header("Chat with your document")
            
            # Display chat history
            for msg in st.session_state.messages:
                with st.chat_message(msg["role"]):
                    st.write(msg["content"])
                    if "images" in msg:
                        for img in msg["images"]:
                            try:
                                image_data = base64.b64decode(img["base64"])
                                st.image(Image.open(io.BytesIO(image_data)))
                            except Exception as e:
                                st.error(f"Error displaying image: {str(e)}")
            
            # Chat input
            if prompt := st.chat_input("Ask about the document"):
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                with st.chat_message("assistant"):
                    texts, images = st.session_state.rag_model.get_relevant_context(prompt)
                    
                    if not texts:
                        response = "I don't have enough context to answer this question."
                    else:
                        prompt_template = ChatPromptTemplate.from_template("""
                        Answer based on this context:
                        {context}
                        
                        Question: {question}
                        """)
                        
                        chain = (
                            RunnableParallel({
                                "context": lambda x: "\n".join(texts),
                                "question": lambda x: prompt
                            })
                            | prompt_template
                            | st.session_state.rag_model.llm
                            | StrOutputParser()
                        )
                        
                        response = chain.invoke({})
                    
                    st.write(response)
                    
                    if images:
                        for img in images:
                            try:
                                image_data = base64.b64decode(img["base64"])
                                st.image(Image.open(io.BytesIO(image_data)))
                            except Exception as e:
                                st.error(f"Error displaying image: {str(e)}")
                    
                    st.session_state.messages.append({
                        "role": "assistant",
                        "content": response,
                        "images": images
                    })
        
        with tab2:
            st.header("Document Summary")
            if st.button("Generate Summary"):
                with st.spinner("Generating summary..."):
                    texts = st.session_state.processed_doc['text_elements']
                    image_summaries = [img['summary'] for img in st.session_state.processed_doc['images']]
                    all_content = texts + image_summaries
                    
                    prompt = ChatPromptTemplate.from_template("""
                    Create a comprehensive medical document summary:
                    
                    Content: {text}
                    
                    Include:
                    1. Key medical findings
                    2. Clinical implications
                    3. Patient information
                    4. Treatment recommendations
                    
                    Format with clear sections and bullet points.
                    """)
                    
                    chain = (
                        RunnableParallel({"text": lambda x: "\n".join(x)})
                        | prompt
                        | st.session_state.rag_model.llm
                        | StrOutputParser()
                    )
                    
                    summary = chain.invoke(all_content)
                    st.markdown(summary)
        
        with tab3:
            st.header("Generate Slides")
            num_slides = st.slider("Number of Slides", 1, 10, 5)
            
            if st.button("Generate Slides"):
                with st.spinner("Generating slides..."):
                    texts = st.session_state.processed_doc['text_elements']
                    image_summaries = [img['summary'] for img in st.session_state.processed_doc['images']]
                    all_content = texts + image_summaries
                    
                    prompt = ChatPromptTemplate.from_template("""
                    Create {num_slides} presentation slides from:
                    
                    {content}
                    
                    Requirements:
                    1. Clear titles
                    2. 3-4 key points per slide
                    3. Medical terminology
                    4. Logical flow
                    
                    Format:
                    Title: [Slide Title]
                    - Point 1: [Content for point 1]
                    - Point 2: [Content for point 2]
                    - Point 3: [Content for point 3]
                    """)
                    
                    chain = (
                        RunnableParallel({
                            "content": lambda x: "\n".join(x['content']),
                            "num_slides": lambda x: x['num_slides']
                        })
                        | prompt
                        | st.session_state.rag_model.llm
                        | StrOutputParser()
                    )
                    
                    result = chain.invoke({
                        "content": all_content,
                        "num_slides": num_slides
                    })
                    st.markdown(result)
        
        with tab4:
            st.header("Document Images")
            if st.session_state.processed_doc['images']:
                for idx, img_info in enumerate(st.session_state.processed_doc['images']):
                    col1, col2 = st.columns([1, 2])
                    with col1:
                        try:
                            image_data = base64.b64decode(img_info['base64'])
                            st.image(Image.open(io.BytesIO(image_data)))
                        except Exception as e:
                            st.error(f"Error displaying image {idx + 1}: {str(e)}")
                    with col2:
                        st.markdown(f"**Image {idx + 1} Analysis:**")
                        st.markdown(img_info.get('summary', 'No analysis available'))
                        if 'metadata' in img_info:
                            st.json(img_info['metadata'])
            else:
                st.info("No images found in the document")
    else:
        st.info("Please upload a document to begin")

if __name__ == "__main__":
    main()